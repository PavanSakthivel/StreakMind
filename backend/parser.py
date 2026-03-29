import io
import PyPDF2
from pptx import Presentation

def extract_text_from_file(file_bytes: bytes, filename: str) -> str:
    ext = filename.split('.')[-1].lower()
    text = ""
    try:
        if ext == 'pdf':
            reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"
        elif ext in ['pptx', 'ppt']:
            prs = Presentation(io.BytesIO(file_bytes))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif ext == 'txt':
            text = file_bytes.decode('utf-8', errors='ignore')
    except Exception as e:
        print(f"Error extracting text from {filename}: {e}")
    return text
